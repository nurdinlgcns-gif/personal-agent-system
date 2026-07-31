from pathlib import Path

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XlsxImage
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Image as PdfImage
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table

ROOT_DIR = Path(__file__).resolve().parents[1]
SAMPLES_DIR = ROOT_DIR / "input" / "samples"

SAMPLES_DIR.mkdir(parents=True, exist_ok=True)


def create_sample_image() -> Path:
    image_path = SAMPLES_DIR / "sample_diagram.png"

    image = Image.new("RGB", (600, 300), color=(245, 247, 250))
    draw = ImageDraw.Draw(image)

    draw.rectangle((40, 40, 560, 260), outline=(30, 80, 160), width=4)
    draw.rectangle((90, 100, 240, 190), outline=(50, 120, 80), width=3)
    draw.rectangle((360, 100, 510, 190), outline=(160, 80, 40), width=3)
    draw.line((240, 145, 360, 145), fill=(30, 30, 30), width=3)

    draw.text((115, 135), "Input", fill=(20, 20, 20))
    draw.text((390, 135), "Output", fill=(20, 20, 20))
    draw.text((225, 40), "Sample Diagram", fill=(20, 20, 20))

    image.save(image_path)

    return image_path


def create_sample_pdf(image_path: Path) -> None:
    output_path = SAMPLES_DIR / "sample_pdf.pdf"
    styles = getSampleStyleSheet()

    document = SimpleDocTemplate(str(output_path), pagesize=A4)
    story = []

    story.append(Paragraph("Technical Manual", styles["Title"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("1. Overview", styles["Heading1"]))
    story.append(
        Paragraph(
            "This PDF contains sample text, a simple table, and a diagram image.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 12))

    table_data = [
        ["Parameter", "Value", "Description"],
        ["WAIT_TIME", "15", "Wait time in seconds"],
        ["SHORT_PATH_RATE", "70", "Short path rate percentage"],
    ]

    story.append(Table(table_data))
    story.append(Spacer(1, 20))

    story.append(Paragraph("2. Diagram", styles["Heading1"]))
    story.append(PdfImage(str(image_path), width=300, height=150))

    document.build(story)


def create_sample_xlsx(image_path: Path) -> None:
    output_path = SAMPLES_DIR / "sample_xlsx.xlsx"

    workbook = Workbook()

    summary_sheet = workbook.active
    summary_sheet.title = "Summary"
    summary_sheet.append(["Item", "Value", "Notes"])
    summary_sheet.append(["Total", 100, "Sample total"])
    summary_sheet.append(["Status", "OK", "Ready for review"])

    detail_sheet = workbook.create_sheet("Detail")
    detail_sheet.append(["ID", "Name", "Qty"])
    detail_sheet.append([1, "Part A", 10])
    detail_sheet.append([2, "Part B", 15])

    image_sheet = workbook.create_sheet("Diagram")
    image_sheet["A1"] = "Embedded diagram sample"

    xlsx_image = XlsxImage(str(image_path))
    xlsx_image.width = 360
    xlsx_image.height = 180
    image_sheet.add_image(xlsx_image, "A3")

    workbook.save(output_path)


def create_sample_pptx(image_path: Path) -> None:
    output_path = SAMPLES_DIR / "sample_pptx.pptx"

    presentation = Presentation()

    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    slide.shapes.title.text = "System Overview"

    content = slide.placeholders[1]
    content.text = "Main process\nDocument conversion\nMarkdown review"

    slide.shapes.add_picture(
        str(image_path),
        PptxInches(5.4),
        PptxInches(1.5),
        width=PptxInches(3),
    )

    notes_text_frame = slide.notes_slide.notes_text_frame
    notes_text_frame.text = "Speaker note: explain the conversion pipeline."

    presentation.save(output_path)


def create_sample_docx(image_path: Path) -> None:
    output_path = SAMPLES_DIR / "sample_docx.docx"

    document = Document()

    document.add_heading("Technical Document", 0)
    document.add_heading("Overview", level=1)
    document.add_paragraph(
        "This DOCX contains heading, paragraph, list, table, and image."
    )

    document.add_paragraph("First bullet", style="List Bullet")
    document.add_paragraph("Second bullet", style="List Bullet")

    table = document.add_table(rows=1, cols=3)
    header = table.rows[0].cells
    header[0].text = "Parameter"
    header[1].text = "Value"
    header[2].text = "Description"

    row = table.add_row().cells
    row[0].text = "WAIT_TIME"
    row[1].text = "15"
    row[2].text = "Wait time in seconds"

    document.add_heading("Diagram", level=1)
    document.add_picture(str(image_path), width=Inches(4))

    document.save(output_path)


def main() -> None:
    image_path = create_sample_image()

    create_sample_pdf(image_path)
    create_sample_xlsx(image_path)
    create_sample_pptx(image_path)
    create_sample_docx(image_path)

    print("Sample files generated:")
    print(f"- {SAMPLES_DIR / 'sample_pdf.pdf'}")
    print(f"- {SAMPLES_DIR / 'sample_xlsx.xlsx'}")
    print(f"- {SAMPLES_DIR / 'sample_pptx.pptx'}")
    print(f"- {SAMPLES_DIR / 'sample_docx.docx'}")


if __name__ == "__main__":
    main()