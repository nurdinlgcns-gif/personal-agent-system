from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from doc_to_md.assets.naming import build_image_name
from doc_to_md.core.models import (
    ContentBlock,
    DocumentSection,
    ExtractedImage,
    ParsedDocument,
)
from doc_to_md.parsers.base_parser import BaseParser


class PptxParser(BaseParser):
    def parse(self, file_path: Path, image_dir: Path) -> ParsedDocument:
        presentation = Presentation(str(file_path))

        parsed = ParsedDocument(
            source_file=file_path.name,
            source_path=file_path,
            source_type="pptx",
            title=file_path.stem,
            sections=[],
            images=[],
            warnings=[
                "PowerPoint SmartArt and complex diagrams may be extracted as images or skipped as structured text.",
                "Shape order may not always match natural human reading order.",
            ],
            metadata={
                "slide_count": len(presentation.slides),
            },
        )

        for slide_index, slide in enumerate(presentation.slides, start=1):
            title = self._get_slide_title(slide) or f"Slide {slide_index}"

            section = DocumentSection(
                heading=f"Slide {slide_index}: {title}",
                level=1,
                blocks=[],
            )

            for shape in slide.shapes:
                if self._is_title_shape(shape, title):
                    continue

                if getattr(shape, "has_text_frame", False):
                    blocks = self._text_frame_to_blocks(shape.text_frame)
                    section.blocks.extend(blocks)

                if getattr(shape, "has_table", False):
                    rows = self._table_to_rows(shape.table)

                    if rows:
                        section.blocks.append(
                            ContentBlock(
                                type="table",
                                content=rows,
                            )
                        )

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = self._extract_picture(
                        shape=shape,
                        image_dir=image_dir,
                        slide_index=slide_index,
                    )

                    if image:
                        parsed.images.append(image)
                        section.blocks.append(
                            ContentBlock(
                                type="image",
                                content={
                                    "alt": image.alt_text,
                                    "path": image.relative_path,
                                },
                            )
                        )

            notes = self._extract_notes(slide)

            if notes:
                section.blocks.append(
                    ContentBlock(
                        type="blockquote",
                        content=f"Speaker Notes:\n{notes}",
                    )
                )

            parsed.sections.append(section)

        return parsed

    def _get_slide_title(self, slide) -> str | None:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()

                if text:
                    first_line = text.splitlines()[0].strip()

                    if first_line:
                        return first_line

        return None

    def _is_title_shape(self, shape, title: str) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False

        text = shape.text.strip()
        return text == title

    def _text_frame_to_blocks(self, text_frame):
        blocks = []

        for paragraph in text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()

            if not text:
                continue

            level = getattr(paragraph, "level", 0) or 0
            indent = "  " * level

            blocks.append(
                ContentBlock(
                    type="bullet_list",
                    content=[f"{indent}{text}"],
                    level=level,
                )
            )

        return blocks

    def _table_to_rows(self, table) -> list[list[str]]:
        rows: list[list[str]] = []

        for row in table.rows:
            values: list[str] = []

            for cell in row.cells:
                values.append(cell.text.strip())

            rows.append(values)

        return rows

    def _extract_picture(
        self,
        shape,
        image_dir: Path,
        slide_index: int,
    ) -> ExtractedImage | None:
        try:
            image = shape.image
            extension = image.ext or "png"

            existing_count = len(
                list(image_dir.glob(f"slide_{slide_index:03d}_image_*"))
            )

            filename = build_image_name(
                prefix=f"slide_{slide_index:03d}",
                index=existing_count + 1,
                extension=extension,
            )

            output_path = image_dir / filename
            output_path.write_bytes(image.blob)

            return ExtractedImage(
                filename=filename,
                relative_path=f"./images/{filename}",
                alt_text=f"Slide {slide_index} image",
            )

        except Exception:
            return None

    def _extract_notes(self, slide) -> str:
        try:
            if not slide.has_notes_slide:
                return ""

            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame

            if not text_frame:
                return ""

            return text_frame.text.strip()

        except Exception:
            return ""
